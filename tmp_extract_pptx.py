import json
from pathlib import Path

pptx_path = Path('public/files/anti-curroption-day.pptx')
out_dir = Path('public/files/anti-curroption-day-assets')
out_dir.mkdir(parents=True, exist_ok=True)

from pptx import Presentation

prs = Presentation(str(pptx_path))
slides_data = []

for i, slide in enumerate(prs.slides, start=1):
    title = None
    texts = []
    images = []

    for shape in slide.shapes:
        if getattr(shape, 'has_text_frame', False) and shape.has_text_frame:
            txt = "\n".join(p.text.strip() for p in shape.text_frame.paragraphs if p.text and p.text.strip())
            if txt:
                texts.append(txt)
                if title is None and len(txt) < 140:
                    title = txt.split('\n')[0]
        if shape.shape_type == 13:  # picture
            image = shape.image
            ext = image.ext or 'png'
            img_name = f"slide-{i:02d}-{len(images)+1}.{ext}"
            img_path = out_dir / img_name
            with open(img_path, 'wb') as f:
                f.write(image.blob)
            images.append('/files/anti-curroption-day-assets/' + img_name)

    slides_data.append({
        'slide': i,
        'title': title or f'Slide {i}',
        'texts': texts,
        'images': images,
    })

json_path = Path('src/pages/anti-curroption-day/content.json')
json_path.parent.mkdir(parents=True, exist_ok=True)
json_path.write_text(json.dumps(slides_data, ensure_ascii=False, indent=2), encoding='utf-8')
print(f'EXTRACTED {len(slides_data)} slides')
print('JSON', json_path)
print('IMG_COUNT', sum(len(s['images']) for s in slides_data))
